import os
import zipfile
from pathlib import Path
import pandas as pd
from docx import Document
from PyPDF2 import PdfReader
import pdfplumber
from pptx import Presentation

# 设置最大文件处理大小 (默认为 100 MB)
# 超过此大小的文件将被跳过，防止内存溢出
MAX_FILE_SIZE_BYTES = 100 * 1024 * 1024

def check_file_size(file_path):
    """
    检查文件大小是否超过限制
    """
    try:
        size = os.path.getsize(file_path)
        if size > MAX_FILE_SIZE_BYTES:
            return False, f"文件过大 ({size / 1024 / 1024:.2f} MB)，超过处理限制 ({MAX_FILE_SIZE_BYTES / 1024 / 1024} MB)"
        return True, ""
    except Exception as e:
        return False, f"无法检查文件大小: {str(e)}"

def read_text_file(file_path):
    """
    读取文本文件（txt, py, md, json, xml, csv等）
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        # 如果UTF-8失败，尝试其他编码
        try:
            with open(file_path, 'r', encoding='gbk') as file:
                return file.read()
        except:
            return f"无法解码文件 {file_path}: 文件编码不支持"
    except Exception as e:
        return f"无法读取文件 {file_path}: {str(e)}"


def read_docx_file(file_path):
    """
    读取Word文档（.docx）
    """
    try:
        doc = Document(file_path)
        content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)
        
        # 添加表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    content.append(" | ".join(row_text))
        
        return "\n".join(content)
    except Exception as e:
        return f"无法读取Word文档 {file_path}: {str(e)}"


def read_pdf_file(file_path):
    """
    读取PDF文件 - 使用pdfplumber作为主要的PDF读取工具，PyPDF2作为备选
    """
    content = []
    pdfplumber_error = None
    
    # 首先尝试使用pdfplumber读取PDF
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                try:
                    text = page.extract_text()
                    if text and text.strip():
                        content.append(f"--- 第{page_num}页 ---")
                        content.append(text.strip())
                except Exception as e:
                    content.append(f"无法提取第{page_num}页内容: {str(e)}")
            
            if content:  # 如果pdfplumber成功读取到内容
                return "\n".join(content)
    except Exception as e:
        pdfplumber_error = e
        pass  # pdfplumber失败，继续尝试其他方法
    
    # 备选方案：使用PyPDF2
    try:
        reader = PdfReader(file_path)
        for page_num, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
                if text and text.strip():
                    content.append(f"--- 第{page_num}页 ---")
                    content.append(text.strip())
            except Exception as e:
                content.append(f"无法提取第{page_num}页内容: {str(e)}")
        
        if content:
            return "\n".join(content)
        else:
            return f"PDF文件 {file_path} 内容为空或无法提取文本"
    except Exception as pypdf2_error:
        if pdfplumber_error:
            return f"无法读取PDF文件 {file_path} (pdfplumber: {str(pdfplumber_error)}, PyPDF2: {str(pypdf2_error)})"
        else:
            return f"无法读取PDF文件 {file_path} (PyPDF2: {str(pypdf2_error)})"


def read_xls_file(file_path):
    """
    读取旧版Excel文件（.xls）
    """
    try:
        # 使用pandas读取.xls文件，header=None 避免将第一行空行误认为表头
        df = pd.read_excel(file_path, sheet_name=None, engine='xlrd', header=None)
        content = []
        
        for sheet_name, sheet_data in df.items():
            # 删除全空的行和列
            sheet_data = sheet_data.dropna(how='all', axis=0)
            sheet_data = sheet_data.dropna(how='all', axis=1)
            
            # 填充 NaN 为空字符串，避免输出 'NaN'
            sheet_data = sheet_data.fillna('')
            
            content.append(f"\n=== 工作表: {sheet_name} ===")
            # 使用 markdown 格式输出，tablefmt="grid" 可以更好地展示表格结构
            try:
                content.append(sheet_data.to_markdown(index=False, header=False, tablefmt="grid"))
            except ImportError:
                # 如果没有安装 tabulate，回退到 to_string
                content.append(sheet_data.to_string(index=False, header=False))
        
        return "\n".join(content)
    except Exception as e:
        return f"无法读取Excel文件 {file_path}: {str(e)}"


def read_excel_file(file_path):
    """
    读取Excel文件（.xlsx, .xls）
    """
    try:
        # 使用pandas读取Excel文件，header=None 避免将第一行空行误认为表头
        df = pd.read_excel(file_path, sheet_name=None, header=None)  # 读取所有工作表
        content = []
        
        for sheet_name, sheet_data in df.items():
            # 删除全空的行和列
            sheet_data = sheet_data.dropna(how='all', axis=0)
            sheet_data = sheet_data.dropna(how='all', axis=1)
            
            # 填充 NaN 为空字符串，避免输出 'NaN'
            sheet_data = sheet_data.fillna('')
            
            content.append(f"\n=== 工作表: {sheet_name} ===")
            # 使用 markdown 格式输出，tablefmt="grid" 可以更好地展示表格结构
            try:
                content.append(sheet_data.to_markdown(index=False, header=False, tablefmt="grid"))
            except ImportError:
                # 如果没有安装 tabulate，回退到 to_string
                content.append(sheet_data.to_string(index=False, header=False))
        
        return "\n".join(content)
    except Exception as e:
        return f"无法读取Excel文件 {file_path}: {str(e)}"


def read_powerpoint_file(file_path):
    """
    读取PowerPoint文件（.pptx）
    """
    try:
        prs = Presentation(file_path)
        content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            content.append(f"\n--- 幻灯片 {slide_num} ---")
            
            # 提取文本框内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
        
        return "\n".join(content)
    except Exception as e:
        return f"无法读取PowerPoint文件 {file_path}: {str(e)}"


def read_csv_file(file_path):
    """
    读取CSV文件
    """
    try:
        df = pd.read_csv(file_path)
        return df.to_string(index=False)
    except Exception as e:
        return f"无法读取CSV文件 {file_path}: {str(e)}"


def read_zip_file(file_path):
    """
    读取Zip文件，尝试识别为Office文档，否则列出内容
    """
    try:
        if not zipfile.is_zipfile(file_path):
             return f"不是有效的Zip文件: {file_path}"
             
        with zipfile.ZipFile(file_path, 'r') as z:
            file_list = z.namelist()
            
            # 检查是否为 Word 文档
            if 'word/document.xml' in file_list:
                return read_docx_file(file_path)
            
            # 检查是否为 Excel 文档
            if 'xl/workbook.xml' in file_list:
                return read_excel_file(file_path)
            
            # 检查是否为 PowerPoint 文档
            if 'ppt/presentation.xml' in file_list:
                return read_powerpoint_file(file_path)
                
            # 如果不是Office文档，列出文件
            content = [f"Zip文件包含 {len(file_list)} 个文件:"]
            for name in file_list:
                # 跳过目录
                if name.endswith('/'):
                    continue
                    
                content.append(f"- {name}")
                
                # 尝试读取文本文件内容
                if name.lower().endswith(('.txt', '.md', '.json', '.xml', '.py', '.csv', '.log')):
                    try:
                        with z.open(name) as f:
                            # 限制读取大小
                            text = f.read(2000).decode('utf-8', errors='ignore')
                            content.append(f"  内容预览:\n{text[:500]}...")
                    except:
                        pass
            
            return "\n".join(content)
    except Exception as e:
        return f"无法读取Zip文件 {file_path}: {str(e)}"


def get_file_content(file_path):
    """
    根据文件类型选择合适的读取方法
    """
    file_path = Path(file_path)
    
    # 检查文件大小
    is_safe, message = check_file_size(file_path)
    if not is_safe:
        return message

    suffix = file_path.suffix.lower()
    
    # 定义支持的文件类型
    file_handlers = {
        # 文本文件
        '.txt': read_text_file,
        '.py': read_text_file,
        '.md': read_text_file,
        '.json': read_text_file,
        '.xml': read_text_file,
        '.html': read_text_file,
        '.css': read_text_file,
        '.js': read_text_file,
        '.log': read_text_file,
        '.ini': read_text_file,
        '.cfg': read_text_file,
        '.conf': read_text_file,
        
        # Word文档
        '.docx': read_docx_file,
        
        # PDF文件
        '.pdf': read_pdf_file,
        
        # Excel文件 - 新版.xlsx
        '.xlsx': read_excel_file,
        # Excel文件 - 旧版.xls
        '.xls': read_xls_file,
        
        # PowerPoint文件
        '.pptx': read_powerpoint_file,
        
        # CSV文件
        '.csv': read_csv_file,
        
        # Zip文件
        '.zip': read_zip_file,
    }
    
    handler = file_handlers.get(suffix)
    if handler:
        return handler(file_path)
    else:
        # 对于不支持的文件类型，尝试作为文本文件读取
        return read_text_file(file_path)

def read_all_files(directory=".", file_extensions=None, exclude_dirs=None):
    """
    读取目录下所有支持的文件内容
    
    参数:
        directory: 要读取的目录路径，默认为当前目录
        file_extensions: 要读取的文件扩展名列表，如 ['.docx', '.pdf']，None表示读取所有支持的文件
        exclude_dirs: 要排除的目录名列表
    
    返回:
        包含所有文件内容的字典，键为文件路径，值为文件内容
    """
    if exclude_dirs is None:
        exclude_dirs = ['.git', '__pycache__', 'node_modules', '.venv', 'venv']
    
    file_contents = {}
    directory_path = Path(directory)
    
    if not directory_path.exists():
        return {"错误": f"目录 {directory} 不存在"}
    
    # 支持的文件扩展名
    supported_extensions = [
        # 文本文件
        '.txt', '.py', '.md', '.json', '.xml', '.html', '.css', '.js', 
        '.log', '.ini', '.cfg', '.conf',
        # Word文档
        '.docx', 
        # PDF文件
        '.pdf',
        # Excel文件 - 新版.xlsx和旧版.xls
        '.xlsx', '.xls',
        # PowerPoint文件
        '.pptx',
        # CSV文件
        '.csv',
        # Zip文件
        '.zip',
    ]
    
    for root, dirs, files in os.walk(directory):
        # 排除指定的目录
        dirs[:] = [d for d in dirs if d not in exclude_dirs]
        
        for file in files:
            file_path = Path(root) / file
            suffix = file_path.suffix.lower()
            
            # 如果指定了文件扩展名，只读取匹配的文件
            if file_extensions:
                if suffix not in file_extensions:
                    continue
            elif suffix not in supported_extensions:
                # 如果没有指定扩展名，只处理支持的文件类型
                continue
            
            # 读取文件内容
            content = get_file_content(file_path)
            file_contents[str(file_path)] = content
    
    return file_contents
