#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建测试文档的统一脚本
包含: .docx, .pdf, .xlsx, .xls, .pptx, .csv, .txt
"""

import os
import pandas as pd
from pathlib import Path
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import xlwt
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.pagesizes import A4

# 尝试注册中文字体，如果没有则使用默认字体（可能不支持中文显示）
try:
    # 尝试常见的Windows中文字体
    pdfmetrics.registerFont(TTFont('SimSun', 'C:/Windows/Fonts/simsun.ttc'))
    FONT_NAME = 'SimSun'
except:
    try:
        pdfmetrics.registerFont(TTFont('SimHei', 'C:/Windows/Fonts/simhei.ttf'))
        FONT_NAME = 'SimHei'
    except:
        print("警告: 未找到常用中文字体，PDF中文可能无法显示")
        FONT_NAME = 'Helvetica'

def create_docx():
    """创建测试Word文档"""
    try:
        doc = Document()
        doc.add_heading('测试Word文档', 0)
        doc.add_paragraph('这是一个测试段落，用于测试Word文档的读取功能。')
        doc.add_paragraph('Python-docx库可以读取Word文档中的文本内容。')
        
        # 添加表格
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = "姓名"
        table.cell(0, 1).text = "年龄"
        table.cell(0, 2).text = "城市"
        table.cell(1, 0).text = "张三"
        table.cell(1, 1).text = "25"
        table.cell(1, 2).text = "北京"
        
        doc.save('测试文档.docx')
        print("✓ 创建测试Word文档: 测试文档.docx")
        return True
    except Exception as e:
        print(f"创建Word文档失败: {e}")
        return False

def create_pdf():
    """创建测试PDF文档 (使用reportlab创建真实PDF)"""
    try:
        c = canvas.Canvas("测试文档.pdf", pagesize=A4)
        c.setFont(FONT_NAME, 16)
        c.drawString(100, 800, "测试PDF文档")
        
        c.setFont(FONT_NAME, 12)
        c.drawString(100, 750, "这是一个测试PDF文档，用于测试PDF文件的读取功能。")
        c.drawString(100, 730, "使用ReportLab库生成。")
        
        # 第二页
        c.showPage()
        c.setFont(FONT_NAME, 12)
        c.drawString(100, 800, "这是第二页的内容。")
        
        c.save()
        print("✓ 创建测试PDF文档: 测试文档.pdf")
        return True
    except Exception as e:
        print(f"创建PDF文档失败: {e}")
        return False

def create_xlsx():
    """创建测试Excel文档 (.xlsx)"""
    try:
        wb = Workbook()
        
        # 第一个工作表
        ws1 = wb.active
        ws1.title = "销售数据"
        ws1['A1'] = "产品"
        ws1['B1'] = "销量"
        ws1['C1'] = "收入"
        ws1['A2'] = "产品A"
        ws1['B2'] = 100
        ws1['C2'] = 5000
        ws1['A3'] = "产品B"
        ws1['B3'] = 150
        ws1['C3'] = 7500
        
        # 第二个工作表
        ws2 = wb.create_sheet("客户信息")
        ws2['A1'] = "客户名称"
        ws2['B1'] = "联系方式"
        ws2['A2'] = "李四"
        ws2['B2'] = "13800138000"
        
        wb.save('测试文档.xlsx')
        print("✓ 创建测试Excel文档: 测试文档.xlsx")
        return True
    except Exception as e:
        print(f"创建Excel(.xlsx)文档失败: {e}")
        return False

def create_xls():
    """创建测试Excel文档 (.xls)"""
    try:
        # 创建工作簿
        workbook = xlwt.Workbook(encoding='utf-8')
        
        # 添加工作表
        worksheet = workbook.add_sheet('测试数据')
        
        # 写入数据
        worksheet.write(0, 0, '产品')
        worksheet.write(0, 1, '销量')
        worksheet.write(0, 2, '收入')
        worksheet.write(1, 0, '产品A')
        worksheet.write(1, 1, 100)
        worksheet.write(1, 2, 5000)
        worksheet.write(2, 0, '产品B')
        worksheet.write(2, 1, 150)
        worksheet.write(2, 2, 7500)
        
        # 添加第二个工作表
        worksheet2 = workbook.add_sheet('客户信息')
        worksheet2.write(0, 0, '客户名称')
        worksheet2.write(0, 1, '联系方式')
        worksheet2.write(1, 0, '李四')
        worksheet2.write(1, 1, '13800138000')
        
        workbook.save('测试文档.xls')
        print("✓ 创建测试Excel文档: 测试文档.xls")
        return True
    except Exception as e:
        print(f"创建Excel(.xls)文档失败: {e}")
        return False

def create_pptx():
    """创建测试PowerPoint文档"""
    try:
        prs = Presentation()
        
        # 第一张幻灯片
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
        title1.text = "测试PowerPoint文档"
        subtitle1.text = "用于测试PPT文件的读取功能"
        
        # 第二张幻灯片
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = slide2.shapes.title
        content2 = slide2.placeholders[1]
        title2.text = "功能介绍"
        content2.text = "python-pptx库可以读取PowerPoint文档中的文本内容"
        
        prs.save('测试文档.pptx')
        print("✓ 创建测试PowerPoint文档: 测试文档.pptx")
        return True
    except Exception as e:
        print(f"创建PowerPoint文档失败: {e}")
        return False

def create_csv():
    """创建测试CSV文档"""
    try:
        data = {
            '姓名': ['张三', '李四', '王五'],
            '年龄': [25, 30, 28],
            '城市': ['北京', '上海', '广州']
        }
        df = pd.DataFrame(data)
        df.to_csv('测试文档.csv', index=False, encoding='utf-8')
        print("✓ 创建测试CSV文档: 测试文档.csv")
        return True
    except Exception as e:
        print(f"创建CSV文档失败: {e}")
        return False

def create_txt():
    """创建测试文本文件"""
    try:
        with open('测试文档.txt', 'w', encoding='utf-8') as f:
            f.write("测试文本文件\n")
            f.write("这是一个测试文本文件，用于测试文本文件的读取功能。\n")
            f.write("Python可以很容易地读取文本文件的内容。\n")
        print("✓ 创建测试文本文件: 测试文档.txt")
        return True
    except Exception as e:
        print(f"创建文本文件失败: {e}")
        return False

def main():
    """主函数：创建所有测试文档"""
    print("开始创建测试文档...")
    
    create_docx()
    create_pdf()
    create_xlsx()
    create_xls()
    create_pptx()
    create_csv()
    create_txt()
    
    print("\n所有测试文档创建完成！")
    print("请运行 main.py 测试文件读取功能。")

if __name__ == "__main__":
    main()
